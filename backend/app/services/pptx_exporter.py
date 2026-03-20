from __future__ import annotations

import re
import struct
import zlib
import hashlib
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.services.template_catalog import resolve_template_assets

TEMPLATE_THEMES = {
    "executive_clean": {
        "bg": RGBColor(245, 248, 252),
        "header": RGBColor(28, 63, 117),
        "header_text": RGBColor(255, 255, 255),
        "card": RGBColor(255, 255, 255),
        "card_alt": RGBColor(237, 244, 255),
        "text": RGBColor(32, 38, 47),
        "muted_text": RGBColor(68, 82, 102),
        "line": RGBColor(210, 220, 235),
        "accent": RGBColor(54, 105, 179),
        "gradient_a": (228, 238, 252),
        "gradient_b": (245, 250, 255),
    },
    "tech_grid": {
        "bg": RGBColor(238, 243, 248),
        "header": RGBColor(21, 48, 75),
        "header_text": RGBColor(243, 248, 255),
        "card": RGBColor(248, 252, 255),
        "card_alt": RGBColor(232, 242, 252),
        "text": RGBColor(23, 37, 52),
        "muted_text": RGBColor(54, 76, 98),
        "line": RGBColor(194, 213, 230),
        "accent": RGBColor(30, 122, 168),
        "gradient_a": (208, 224, 242),
        "gradient_b": (238, 245, 252),
    },
    "dark_focus": {
        "bg": RGBColor(30, 35, 44),
        "header": RGBColor(41, 54, 78),
        "header_text": RGBColor(245, 248, 252),
        "card": RGBColor(49, 58, 73),
        "card_alt": RGBColor(57, 69, 87),
        "text": RGBColor(241, 245, 252),
        "muted_text": RGBColor(213, 224, 239),
        "line": RGBColor(95, 114, 141),
        "accent": RGBColor(78, 195, 255),
        "gradient_a": (22, 26, 34),
        "gradient_b": (42, 50, 64),
    },
    "warm_report": {
        "bg": RGBColor(248, 243, 236),
        "header": RGBColor(114, 84, 61),
        "header_text": RGBColor(253, 249, 243),
        "card": RGBColor(255, 251, 246),
        "card_alt": RGBColor(249, 238, 226),
        "text": RGBColor(64, 51, 42),
        "muted_text": RGBColor(92, 73, 58),
        "line": RGBColor(222, 201, 182),
        "accent": RGBColor(193, 120, 69),
        "gradient_a": (245, 229, 209),
        "gradient_b": (252, 245, 234),
    },
    "ocean_blue": {
        "bg": RGBColor(237, 244, 255),
        "header": RGBColor(19, 59, 115),
        "header_text": RGBColor(247, 251, 255),
        "card": RGBColor(250, 253, 255),
        "card_alt": RGBColor(232, 242, 255),
        "text": RGBColor(22, 43, 69),
        "muted_text": RGBColor(55, 83, 117),
        "line": RGBColor(193, 212, 236),
        "accent": RGBColor(47, 125, 209),
        "gradient_a": (220, 235, 255),
        "gradient_b": (241, 248, 255),
    },
    "forest_brief": {
        "bg": RGBColor(238, 247, 241),
        "header": RGBColor(31, 79, 59),
        "header_text": RGBColor(244, 252, 247),
        "card": RGBColor(250, 255, 252),
        "card_alt": RGBColor(233, 246, 238),
        "text": RGBColor(31, 56, 45),
        "muted_text": RGBColor(66, 98, 83),
        "line": RGBColor(191, 219, 205),
        "accent": RGBColor(47, 141, 101),
        "gradient_a": (221, 241, 230),
        "gradient_b": (243, 252, 247),
    },
    "sunset_orange": {
        "bg": RGBColor(255, 244, 234),
        "header": RGBColor(106, 55, 25),
        "header_text": RGBColor(255, 250, 245),
        "card": RGBColor(255, 252, 247),
        "card_alt": RGBColor(255, 240, 225),
        "text": RGBColor(79, 45, 24),
        "muted_text": RGBColor(116, 73, 45),
        "line": RGBColor(236, 204, 176),
        "accent": RGBColor(224, 122, 47),
        "gradient_a": (255, 231, 207),
        "gradient_b": (255, 247, 236),
    },
    "internet_ops": {
        "bg": RGBColor(238, 244, 255),
        "header": RGBColor(28, 58, 116),
        "header_text": RGBColor(247, 251, 255),
        "card": RGBColor(250, 253, 255),
        "card_alt": RGBColor(232, 240, 255),
        "text": RGBColor(24, 45, 74),
        "muted_text": RGBColor(66, 91, 128),
        "line": RGBColor(195, 212, 239),
        "accent": RGBColor(63, 111, 216),
        "gradient_a": (219, 232, 255),
        "gradient_b": (242, 247, 255),
    },
    "manufacturing_ops": {
        "bg": RGBColor(238, 247, 245),
        "header": RGBColor(25, 77, 68),
        "header_text": RGBColor(244, 252, 250),
        "card": RGBColor(250, 255, 254),
        "card_alt": RGBColor(231, 245, 241),
        "text": RGBColor(27, 61, 55),
        "muted_text": RGBColor(67, 104, 97),
        "line": RGBColor(191, 219, 212),
        "accent": RGBColor(46, 155, 134),
        "gradient_a": (220, 241, 236),
        "gradient_b": (242, 251, 249),
    },
    "investor_pitch": {
        "bg": RGBColor(255, 243, 234),
        "header": RGBColor(110, 52, 23),
        "header_text": RGBColor(255, 250, 245),
        "card": RGBColor(255, 252, 247),
        "card_alt": RGBColor(255, 239, 226),
        "text": RGBColor(82, 44, 22),
        "muted_text": RGBColor(122, 78, 50),
        "line": RGBColor(238, 204, 178),
        "accent": RGBColor(240, 138, 60),
        "gradient_a": (255, 229, 207),
        "gradient_b": (255, 247, 237),
    },
}

ICON_STYLES = ["circle", "square", "diamond", "arrow"]
DEFENSE_COVER_SUBTITLE = "\u7855\u58eb\u5b66\u4f4d\u8bba\u6587\u7b54\u8fa9"


def _derive_theme(template_id: str) -> Dict:
    digest = hashlib.md5(template_id.encode("utf-8")).digest()

    def _pick(seed: int, low: int, span: int) -> int:
        return low + (digest[seed] % span)

    bg_rgb = (_pick(0, 228, 24), _pick(1, 233, 20), _pick(2, 238, 17))
    header_rgb = (_pick(3, 18, 44), _pick(4, 45, 40), _pick(5, 70, 35))
    text_rgb = (_pick(6, 24, 40), _pick(7, 38, 34), _pick(8, 52, 30))
    muted_rgb = (
        min(235, text_rgb[0] + 28),
        min(235, text_rgb[1] + 24),
        min(235, text_rgb[2] + 22),
    )
    line_rgb = (_pick(9, 188, 30), _pick(10, 206, 24), _pick(11, 217, 22))
    accent_rgb = (_pick(12, 45, 130), _pick(13, 95, 120), _pick(14, 120, 115))

    return {
        "bg": RGBColor(*bg_rgb),
        "header": RGBColor(*header_rgb),
        "header_text": RGBColor(248, 251, 255),
        "card": RGBColor(252, 254, 255),
        "card_alt": RGBColor(max(225, bg_rgb[0] - 10), max(228, bg_rgb[1] - 8), max(232, bg_rgb[2] - 6)),
        "text": RGBColor(*text_rgb),
        "muted_text": RGBColor(*muted_rgb),
        "line": RGBColor(*line_rgb),
        "accent": RGBColor(*accent_rgb),
        "gradient_a": (max(210, bg_rgb[0] - 14), max(218, bg_rgb[1] - 12), max(224, bg_rgb[2] - 10)),
        "gradient_b": (min(252, bg_rgb[0] + 10), min(254, bg_rgb[1] + 8), min(255, bg_rgb[2] + 7)),
    }


def _theme(template_id: str) -> Dict:
    return TEMPLATE_THEMES.get(template_id, _derive_theme(template_id))

def _truncate(text: str, limit: int = 100) -> str:
    t = (text or "").strip()
    return t if len(t) <= limit else t[: limit - 3].rstrip() + "..."


def _normalize_toc_item(raw: str) -> str:
    txt = str(raw or "").strip()
    if not txt:
        return txt
    txt = re.sub(r"^\u7b2c\s*\d+\s*\u9875[\uff1a:]\s*", "", txt)
    txt = re.sub(r"^\d+\s*[\.\u3001\)\uff09]\s*", "", txt)
    return txt.strip()

def _png_chunk(chunk_type: bytes, data: bytes) -> bytes:
    crc = zlib.crc32(chunk_type + data) & 0xFFFFFFFF
    return struct.pack("!I", len(data)) + chunk_type + data + struct.pack("!I", crc)


def _write_gradient_png(path: Path, width: int, height: int, a: Tuple[int, int, int], b: Tuple[int, int, int]) -> None:
    rows = []
    for y in range(height):
        t = y / max(1, height - 1)
        r0 = int(a[0] * (1 - t) + b[0] * t)
        g0 = int(a[1] * (1 - t) + b[1] * t)
        b0 = int(a[2] * (1 - t) + b[2] * t)
        row = bytearray([0])
        for x in range(width):
            wave = ((x * 11 + y * 7) % 29) - 14
            row.extend((
                max(0, min(255, r0 + wave // 3)),
                max(0, min(255, g0 + wave // 4)),
                max(0, min(255, b0 + wave // 5)),
            ))
        rows.append(bytes(row))

    raw = b"".join(rows)
    compressed = zlib.compress(raw, 9)
    signature = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack("!IIBBBBB", width, height, 8, 2, 0, 0, 0)
    png = signature + _png_chunk(b"IHDR", ihdr) + _png_chunk(b"IDAT", compressed) + _png_chunk(b"IEND", b"")
    path.write_bytes(png)


def _ensure_template_background(template_id: str, theme: Dict) -> Path:
    out_dir = Path(__file__).resolve().parents[2] / "assets" / "template_bgs"
    out_dir.mkdir(parents=True, exist_ok=True)
    img = out_dir / f"{template_id}.png"
    if not img.exists():
        _write_gradient_png(img, 1600, 900, theme["gradient_a"], theme["gradient_b"])
    return img


def _add_bg_image(slide, bg_path: Optional[Path]) -> None:
    if not bg_path or not bg_path.exists():
        return
    slide.shapes.add_picture(str(bg_path), Inches(0), Inches(0), Inches(13.33), Inches(7.5))

def _set_card_style(shape, theme: Dict, alt: bool = False) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["card_alt" if alt else "card"]
    shape.line.color.rgb = theme["line"]
    shape.line.width = Pt(1.2)


def _add_title_bar(slide, title: str, theme: Dict) -> None:
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = theme["header"]
    header.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.75), Inches(0.42), Inches(11.8), Inches(0.54))
    tf = box.text_frame
    tf.clear()
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = _truncate(title, 78)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = theme["header_text"]


def _make_slide(prs: Presentation, title: str, theme: Dict, bg_path: Optional[Path]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_image(slide, bg_path)
    _add_title_bar(slide, title, theme)
    return slide


def _extract_topic(raw_title: str) -> str:
    return raw_title.split(" - ", 1)[0].strip() if " - " in raw_title else raw_title.strip()


def _render_cover(prs: Presentation, topic: str, subtitle: str, theme: Dict, bg_path: Optional[Path]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_image(slide, bg_path)

    veil = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    veil.fill.solid()
    veil.fill.fore_color.rgb = theme["header"]
    veil.fill.transparency = 0.32
    veil.line.fill.background()

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.1), Inches(11.6), Inches(5.2))
    _set_card_style(panel, theme, alt=True)

    box = slide.shapes.add_textbox(Inches(1.35), Inches(2.0), Inches(10.7), Inches(2.5))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = _truncate(topic, 60)
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = theme["text"]
    p0.alignment = PP_ALIGN.LEFT

    p1 = tf.add_paragraph()
    p1.text = _truncate(subtitle, 80)
    p1.font.size = Pt(20)
    p1.font.color.rgb = theme["muted_text"]
    p1.space_before = Pt(12)


def _render_toc(prs: Presentation, topic: str, outline: List[str], theme: Dict, bg_path: Optional[Path]):
    slide = _make_slide(prs, f"{topic} | Agenda", theme, bg_path)
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.8), Inches(5.7))
    _set_card_style(card, theme)

    items = [x for x in outline if x and ("cover" not in x.lower()) and ("agenda" not in x.lower())]
    items = items[:8]
    link_shapes = []

    y = 1.8
    for idx, item in enumerate(items, start=1):
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.15), Inches(y + 0.05), Inches(0.35), Inches(0.35))
        icon.fill.solid()
        icon.fill.fore_color.rgb = theme["accent"]
        icon.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(1.22), Inches(y + 0.07), Inches(0.2), Inches(0.2))
        p = num_box.text_frame.paragraphs[0]
        p.text = str(idx)
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        txt = slide.shapes.add_textbox(Inches(1.6), Inches(y), Inches(10.2), Inches(0.52))
        tp = txt.text_frame.paragraphs[0]
        tp.text = _truncate(_normalize_toc_item(item), 68)
        tp.font.size = Pt(19)
        tp.font.color.rgb = theme["text"]
        link_shapes.append(txt)
        y += 0.63

    return link_shapes


def _icon_shape(slide, style: str, x: float, y: float, color: RGBColor):
    st = {
        "circle": MSO_AUTO_SHAPE_TYPE.OVAL,
        "square": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
        "arrow": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
    }.get(style, MSO_AUTO_SHAPE_TYPE.OVAL)
    icon = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(0.22), Inches(0.22))
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()


def _add_icon_bullets(slide, points: List[str], theme: Dict, box: Tuple[float, float, float, float], start_style: int = 0):
    x, y, w, h = box
    clean = [_truncate(p, 88) for p in points if p] or ["TBD point"]
    clean = clean[:5]
    line_h = min(0.95, max(0.62, (h - 0.2) / len(clean)))

    cy = y
    for idx, text in enumerate(clean):
        style = ICON_STYLES[(start_style + idx) % len(ICON_STYLES)]
        _icon_shape(slide, style, x, cy + 0.12, theme["accent"])
        tbox = slide.shapes.add_textbox(Inches(x + 0.33), Inches(cy), Inches(w - 0.35), Inches(line_h))
        tf = tbox.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(17)
        p.font.color.rgb = theme["muted_text"]
        p.alignment = PP_ALIGN.LEFT
        cy += line_h


def _render_summary(prs: Presentation, slide_data: Dict, theme: Dict, bg_path: Optional[Path]):
    slide = _make_slide(prs, slide_data.get("title", "Content"), theme, bg_path)
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.5), Inches(11.6), Inches(5.75))
    _set_card_style(card, theme)
    _add_icon_bullets(slide, slide_data.get("bullets", []), theme, (1.2, 1.95, 10.9, 4.9), start_style=0)
    return slide


def _render_risk(prs: Presentation, slide_data: Dict, theme: Dict, bg_path: Optional[Path]):
    slide = _make_slide(prs, slide_data.get("title", "Risk"), theme, bg_path)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.7))
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.7))
    _set_card_style(left, theme)
    _set_card_style(right, theme, alt=True)

    lt = slide.shapes.add_textbox(Inches(1.1), Inches(1.75), Inches(4.8), Inches(0.45))
    lp = lt.text_frame.paragraphs[0]
    lp.text = "Top Risks"
    lp.font.size = Pt(20)
    lp.font.bold = True
    lp.font.color.rgb = theme["text"]

    rt = slide.shapes.add_textbox(Inches(7.1), Inches(1.75), Inches(4.8), Inches(0.45))
    rp = rt.text_frame.paragraphs[0]
    rp.text = "Mitigation"
    rp.font.size = Pt(20)
    rp.font.bold = True
    rp.font.color.rgb = theme["text"]

    _add_icon_bullets(slide, slide_data.get("bullets", []), theme, (1.1, 2.25, 5.0, 4.6), start_style=1)
    right_points = slide_data.get("evidence", []) or slide_data.get("bullets", [])
    _add_icon_bullets(slide, right_points, theme, (7.1, 2.25, 5.0, 4.6), start_style=2)
    return slide


def _render_timeline(prs: Presentation, slide_data: Dict, theme: Dict, bg_path: Optional[Path]):
    slide = _make_slide(prs, slide_data.get("title", "Timeline"), theme, bg_path)
    track = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.1), Inches(3.62), Inches(11.0), Inches(0.08))
    track.fill.solid()
    track.fill.fore_color.rgb = theme["line"]
    track.line.fill.background()

    points = [_truncate(x, 54) for x in slide_data.get("bullets", [])[:4]] or ["Stage detail TBD"]
    n = len(points)
    for idx, text in enumerate(points):
        x = 1.2 + idx * (10.6 / max(1, n - 1))
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(3.42), Inches(0.42), Inches(0.42))
        node.fill.solid()
        node.fill.fore_color.rgb = theme["accent"]
        node.line.fill.background()

        top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x - 0.5), Inches(1.95), Inches(1.5), Inches(1.2))
        _set_card_style(top, theme, alt=(idx % 2 == 1))
        tf = top.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = f"Stage {idx + 1}"
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = theme["muted_text"]
        p1 = tf.add_paragraph()
        p1.text = text
        p1.font.size = Pt(12)
        p1.font.color.rgb = theme["text"]
    return slide


def _chart_payload(slide_data: Dict) -> Tuple[List[str], List[float], str]:
    chart = slide_data.get("chart_data")
    if isinstance(chart, dict):
        labels = [str(x).strip() for x in chart.get("labels", []) if str(x).strip()]
        values: List[float] = []
        for v in chart.get("values", []):
            try:
                values.append(float(v))
            except (TypeError, ValueError):
                continue
        unit = str(chart.get("unit", "") or "")
        if labels and values and len(labels) == len(values):
            return labels[:6], values[:6], unit

    labels = ["Metric 1", "Metric 2", "Metric 3"]
    values = [60.0, 72.0, 84.0]
    return labels, values, ""


def _render_data(prs: Presentation, slide_data: Dict, theme: Dict, bg_path: Optional[Path]):
    slide = _make_slide(prs, slide_data.get("title", "Data"), theme, bg_path)

    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.3), Inches(5.7))
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(1.5), Inches(6.0), Inches(5.7))
    _set_card_style(left, theme)
    _set_card_style(right, theme, alt=True)

    _add_icon_bullets(slide, slide_data.get("bullets", []), theme, (1.1, 1.95, 4.7, 4.9), start_style=0)

    labels, values, unit = _chart_payload(slide_data)
    max_value = max(values) if values else 1.0
    max_value = max(max_value, 1.0)

    for idx, (label, value) in enumerate(zip(labels, values)):
        y = 2.15 + idx * 0.72
        label_box = slide.shapes.add_textbox(Inches(6.85), Inches(y), Inches(1.45), Inches(0.55))
        lp = label_box.text_frame.paragraphs[0]
        lp.text = _truncate(label, 16)
        lp.font.size = Pt(12)
        lp.font.color.rgb = theme["muted_text"]

        width = max(0.3, (value / max_value) * 4.9)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(8.3),
            Inches(y + 0.06),
            Inches(width),
            Inches(0.42),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme["accent"]
        bar.line.fill.background()

        val_box = slide.shapes.add_textbox(Inches(8.4 + width), Inches(y), Inches(1.6), Inches(0.55))
        vp = val_box.text_frame.paragraphs[0]
        vp.text = f"{value:.1f}{unit}"
        vp.font.size = Pt(12)
        vp.font.bold = True
        vp.font.color.rgb = theme["text"]

    return slide


def _content_slides(slides: List[Dict]) -> List[Dict]:
    out = []
    for s in slides:
        title = str(s.get("title", ""))
        if s.get("slide_type") == "title":
            continue
        low = title.lower()
        if "cover" in low or "agenda" in low:
            continue
        out.append(s)
    return out



def _text_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def _shape_area(shape) -> int:
    return int(getattr(shape, "width", 0)) * int(getattr(shape, "height", 0))


def _slide_text_capacity(slide) -> int:
    score = 0
    for shp in _text_shapes(slide):
        score += 2_000_000 if getattr(shp, "is_placeholder", False) else 0
        score += _shape_area(shp)
    return score


def _pick_title_shape(slide):
    title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
    text_shapes = _text_shapes(slide)

    for shape in text_shapes:
        if getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.type in title_types:
                    return shape
            except Exception:
                pass

    if not text_shapes:
        return None

    ordered = sorted(text_shapes, key=lambda s: (int(getattr(s, "top", 0)), -int(getattr(s, "width", 0))))
    return ordered[0]


def _pick_body_shapes(slide, title_shape=None):
    body_type_names = (
        "BODY",
        "OBJECT",
        "SUBTITLE",
        "VERTICAL_OBJECT",
        "VERTICAL_BODY",
        "VERTICAL_TITLE_AND_TEXT",
        "CONTENT",
    )
    body_types = {getattr(PP_PLACEHOLDER, name) for name in body_type_names if hasattr(PP_PLACEHOLDER, name)}
    text_shapes = [s for s in _text_shapes(slide) if s is not title_shape]

    placeholder_bodies = []
    for shape in text_shapes:
        if getattr(shape, "is_placeholder", False):
            try:
                if not body_types or shape.placeholder_format.type in body_types:
                    placeholder_bodies.append(shape)
            except Exception:
                pass

    if placeholder_bodies:
        return sorted(placeholder_bodies, key=_shape_area, reverse=True)
    return sorted(text_shapes, key=_shape_area, reverse=True)


def _set_shape_text(shape, text: str) -> None:
    if not shape:
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    tf.paragraphs[0].text = _truncate(text, 96)


def _set_shape_bullets(shape, lines: List[str], limit: int = 6) -> None:
    if not shape:
        return
    clean = [_truncate(x, 68) for x in lines if x]
    if not clean:
        clean = [" "]
    clean = clean[:limit]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    tf.paragraphs[0].text = clean[0]
    for item in clean[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0


def _is_default_placeholder_text(text: str) -> bool:
    t = (text or "").strip().lower()
    if not t:
        return True
    tokens = (
        "\u70b9\u51fb\u6dfb\u52a0\u6807\u9898",
        "\u70b9\u51fb\u6dfb\u52a0\u6587\u672c",
        "\u70b9\u51fb\u6dfb\u52a0\u5185\u5bb9",
        "\u60a8\u7684\u5185\u5bb9\u6253\u5728\u8fd9\u91cc",
        "add title",
        "add text",
        "click to add",
        "or paste your text",
        "content here",
    )
    if any(x in t for x in tokens):
        return True

    # Common template artifacts left in unused shapes.
    lines = [ln.strip().lower() for ln in re.split(r"[\r\n]+", t) if ln.strip()]
    if lines and len(lines) <= 4:
        short_token = re.compile(r"^(?:[a-z]|[ivx]{1,4}|\d{1,2}|\d{6}|\d{4}\.\d{2})$")
        if all(short_token.match(ln) for ln in lines):
            return True

    return False

def _clear_shape_text(shape) -> None:
    if not shape or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True


def _cleanup_unused_text_shapes(slide, keep_shapes) -> None:
    keep_ids = {id(s) for s in keep_shapes if s is not None}
    for shp in _text_shapes(slide):
        if id(shp) in keep_ids:
            continue
        if getattr(shp, "is_placeholder", False):
            _clear_shape_text(shp)
            continue
        txt = ""
        try:
            txt = shp.text_frame.text
        except Exception:
            txt = ""
        if _is_default_placeholder_text(txt):
            _clear_shape_text(shp)


def _fill_lines_across_shapes(shapes: List, lines: List[str], max_per_shape: int = 3) -> List:
    used = []
    clean = [x for x in lines if x]
    if not shapes:
        return used
    if not clean:
        _set_shape_bullets(shapes[0], [" "])
        return [shapes[0]]

    # If template provides enough boxes, map one bullet to one box for better layout fidelity.
    if len(shapes) >= len(clean):
        for shp, item in zip(shapes, clean):
            _set_shape_bullets(shp, [item], limit=1)
            used.append(shp)
        return used

    idx = 0
    for shp in shapes:
        chunk = clean[idx: idx + max_per_shape]
        if not chunk:
            break
        _set_shape_bullets(shp, chunk, limit=max_per_shape)
        used.append(shp)
        idx += max_per_shape

    return used

def _delete_slide(prs: Presentation, idx: int) -> None:
    slide_id = prs.slides._sldIdLst[idx]
    rid = slide_id.rId
    prs.part.drop_rel(rid)
    del prs.slides._sldIdLst[idx]


def _template_slide_text(slide) -> str:
    chunks: List[str] = []
    for shp in _text_shapes(slide):
        try:
            txt = (shp.text_frame.text or "").strip()
        except Exception:
            txt = ""
        if txt:
            chunks.append(txt)
    return " ".join(chunks).lower()


def _template_slide_role(slide, idx: int) -> str:
    txt = _template_slide_text(slide)
    if any(k in txt for k in ("agenda", "contents", "??")):
        return "toc"
    if any(k in txt for k in ("??", "??", "thesis", "??")) and idx <= 1:
        return "cover"
    if any(k in txt for k in ("timeline", "???", "??", "??", "roadmap")):
        return "timeline"
    if any(k in txt for k in ("risk", "??", "??", "??", "mitigation")):
        return "risk"
    if any(k in txt for k in ("data", "??", "??", "??", "??", "%", "??", "analysis")):
        return "data"
    return "summary"


def _template_body_slot_count(slide) -> int:
    title_shape = _pick_title_shape(slide)
    return len(_pick_body_shapes(slide, title_shape))


def _pick_cover_and_toc(descriptors: List[Dict]) -> tuple[int, int]:
    def cover_score(d: Dict) -> int:
        txt = d["text"]
        score = (4 if d["idx"] == 0 else 0) + (2 if d["has_title"] else 0)
        if any(k in txt for k in ("??", "??", "thesis", "??")):
            score += 3
        if d["role"] == "cover":
            score += 2
        return score

    def toc_score(d: Dict) -> int:
        txt = d["text"]
        score = (3 if d["idx"] == 1 else 0) + (1 if d["has_title"] else 0)
        if any(k in txt for k in ("??", "agenda", "contents")):
            score += 4
        if d["role"] == "toc":
            score += 2
        return score

    cover = max(descriptors, key=cover_score)["idx"] if descriptors else 0
    remain = [d for d in descriptors if d["idx"] != cover]
    toc = (max(remain, key=toc_score)["idx"] if remain else cover)
    return cover, toc


def _content_match_score(d: Dict, slide_type: str) -> int:
    role = d["role"]
    score = 0
    if slide_type == "timeline":
        score += 8 if role == "timeline" else 0
    elif slide_type == "data":
        score += 8 if role == "data" else 0
    elif slide_type == "risk":
        score += 8 if role == "risk" else 0
    else:
        score += 5 if role == "summary" else 0

    # Secondary compatibility.
    if role in {"summary", "data", "risk", "timeline"}:
        score += 2

    score += min(4, d["body_slots"])
    score += min(4, d["capacity"] // 2_500_000)
    return score


def _select_template_slide_indices(prs: Presentation, body_slides: List[Dict]) -> List[int]:
    total_needed = 2 + len(body_slides)
    total = len(prs.slides)
    if total <= total_needed:
        return list(range(total))

    descriptors: List[Dict] = []
    for idx, slide in enumerate(prs.slides):
        title_shape = _pick_title_shape(slide)
        descriptors.append(
            {
                "idx": idx,
                "role": _template_slide_role(slide, idx),
                "text": _template_slide_text(slide),
                "capacity": _slide_text_capacity(slide),
                "has_title": bool(title_shape),
                "body_slots": _template_body_slot_count(slide),
            }
        )

    cover_idx, toc_idx = _pick_cover_and_toc(descriptors)
    selected = [cover_idx]
    used = {cover_idx}
    if toc_idx not in used:
        selected.append(toc_idx)
        used.add(toc_idx)

    cursor = max(selected) if selected else -1

    for payload in body_slides:
        slide_type = str(payload.get("slide_type", "summary") or "summary")
        candidates = [d for d in descriptors if d["idx"] not in used and d["idx"] > cursor]
        if not candidates:
            candidates = [d for d in descriptors if d["idx"] not in used]
        if not candidates:
            break

        pick = max(candidates, key=lambda d: _content_match_score(d, slide_type))
        selected.append(pick["idx"])
        used.add(pick["idx"])
        cursor = pick["idx"]

    if len(selected) < total_needed:
        rest = [d["idx"] for d in sorted(descriptors, key=lambda d: d["capacity"], reverse=True) if d["idx"] not in used]
        for idx in rest:
            selected.append(idx)
            used.add(idx)
            if len(selected) >= total_needed:
                break

    return selected[:total_needed]

def _export_with_template_pages(
    slides: List[Dict],
    out_path: Path,
    template_pptx_path: Path,
    topic: str = "",
    outline: Optional[List[str]] = None,
) -> str:
    prs = Presentation(str(template_pptx_path))
    body_slides = _content_slides(slides)
    toc_source = outline[:] if outline else [s.get("title", "") for s in body_slides]
    total_needed = 2 + len(body_slides)

    while len(prs.slides) < total_needed:
        layout_idx = 0 if len(prs.slide_layouts) > 0 else 6
        prs.slides.add_slide(prs.slide_layouts[layout_idx])

    keep_indices = _select_template_slide_indices(prs, body_slides)
    remove_indices = [i for i in range(len(prs.slides)) if i not in set(keep_indices)]
    for idx in sorted(remove_indices, reverse=True):
        _delete_slide(prs, idx)

    base_title = topic.strip() if topic else (_extract_topic(body_slides[0].get("title", "Report")) if body_slides else "Report")

    cover = prs.slides[0]
    cover_title = _pick_title_shape(cover)
    cover_bodies = _pick_body_shapes(cover, cover_title)
    used_cover = [cover_title]
    _set_shape_text(cover_title, base_title)
    if cover_bodies:
        _set_shape_text(cover_bodies[0], DEFENSE_COVER_SUBTITLE)
        used_cover.append(cover_bodies[0])
    _cleanup_unused_text_shapes(cover, used_cover)
    if cover.has_notes_slide:
        cover.notes_slide.notes_text_frame.text = ""

    if len(prs.slides) > 1:
        toc = prs.slides[1]
        toc_title = _pick_title_shape(toc)
        toc_bodies = _pick_body_shapes(toc, toc_title)
        _set_shape_text(toc_title, "Agenda")
        used = [toc_title]
        if toc_bodies:
            toc_lines = [f"{idx + 1}. {_normalize_toc_item(str(item))}" for idx, item in enumerate(toc_source[:10])]
            used.extend(_fill_lines_across_shapes(toc_bodies, toc_lines, max_per_shape=2))
        _cleanup_unused_text_shapes(toc, used)
        if toc.has_notes_slide:
            toc.notes_slide.notes_text_frame.text = ""

    for idx, payload in enumerate(body_slides, start=2):
        if idx >= len(prs.slides):
            break
        slide = prs.slides[idx]
        title_shape = _pick_title_shape(slide)
        body_shapes = _pick_body_shapes(slide, title_shape)
        used = [title_shape]
        _set_shape_text(title_shape, payload.get("title", ""))
        if body_shapes:
            used.extend(_fill_lines_across_shapes(body_shapes, payload.get("bullets", []), max_per_shape=3))
        _cleanup_unused_text_shapes(slide, used)
        if slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = payload.get("notes", "")

    prs.save(str(out_path))
    return out_path.name
def export_slides_to_pptx(
    slides: List[Dict],
    out_path: Path,
    template_id: str = "executive_clean",
    topic: str = "",
    outline: List[str] | None = None,
) -> str:
    assets = resolve_template_assets(template_id)
    template_pptx_path = assets.get("pptx_path")
    custom_bg_path = assets.get("bg_path")

    if template_pptx_path and template_pptx_path.exists():
        return _export_with_template_pages(
            slides=slides,
            out_path=out_path,
            template_pptx_path=template_pptx_path,
            topic=topic,
            outline=outline,
        )

    prs = Presentation()

    theme = _theme(template_id)
    bg_path = custom_bg_path if custom_bg_path and custom_bg_path.exists() else _ensure_template_background(template_id, theme)

    body_slides = _content_slides(slides)
    base_title = topic.strip() if topic else (_extract_topic(body_slides[0].get("title", "Report")) if body_slides else "Report")

    _render_cover(prs, base_title, DEFENSE_COVER_SUBTITLE, theme, bg_path)

    toc_source = outline[:] if outline else [s.get("title", "") for s in body_slides]
    toc_link_shapes = _render_toc(prs, base_title, toc_source, theme, bg_path)

    content_rendered_slides = []
    for slide_data in body_slides:
        slide_type = slide_data.get("slide_type", "summary")
        if slide_type == "risk":
            slide = _render_risk(prs, slide_data, theme, bg_path)
        elif slide_type == "data":
            slide = _render_data(prs, slide_data, theme, bg_path)
        elif slide_type in ("timeline", "status"):
            slide = _render_timeline(prs, slide_data, theme, bg_path)
        else:
            slide = _render_summary(prs, slide_data, theme, bg_path)

        content_rendered_slides.append(slide)
        if slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")

    for idx, shape in enumerate(toc_link_shapes):
        if idx >= len(content_rendered_slides):
            break
        shape.click_action.target_slide = content_rendered_slides[idx]

    prs.save(str(out_path))
    return out_path.name












